#!/usr/bin/env python3
"""
Générateur d'affiches anniversaire UrbanSoccer.

Lit la feuille de route Excel et génère les affiches PPTX + PDF
pour chaque anniversaire du jour.

Usage :
    python generate_birthday_posters.py                        # aujourd'hui
    python generate_birthday_posters.py --date 22.02.26        # date spécifique
    python generate_birthday_posters.py --sheet "15.02.26 (2)" # onglet spécifique

Le script :
  1. Lit l'onglet correspondant dans le fichier Excel « feuille de route »
  2. Génère un PPTX avec une affiche par enfant (template adapté à la formule)
  3. Convertit en PDF via LibreOffice (si disponible)
"""

import argparse
import copy
import io
import os
import re
import subprocess
import sys
from datetime import datetime

import openpyxl
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn
from pptx.parts.slide import SlidePart

# ── Chemins par défaut ────────────────────────────────────────────────────────

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
EXCEL_PATH = os.path.join(SCRIPT_DIR, "Anniversaires",
                          "Anniversaire - Feuille de route (1).xlsx")
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "Anniversaires",
                             "Affiche-Bienvenue-Anniversaire-2022 (1) 3-5.pptx")
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "Anniversaires", "output")

# ── Template slide indices (0-based) dans le fichier PPTX ─────────────────────
#
# Orange (Ligue 1 / Champions League) : slides 4-17 (4 shapes)
#   → shape[3] = prénom, 96pt, couleur FE7831
#
# FFF bienvenue : slides 18, 20 (3 shapes)
#   → shape[2] = prénom, font Arsilon, 115pt, couleur FFFFFF
#
# FFF certificat : slides 19, 21 (2 shapes)
#   → shape[1] = code supporter, 40pt, couleur thème BACKGROUND_1
#
# Bump : slides 0-2 (9 shapes, double page + branding Bump)
#   → shape[6] = prénom (page visible), 96pt, couleur FE7831

TEMPLATE_ORANGE = 4
TEMPLATE_FFF_BIENVENUE = 18
TEMPLATE_FFF_CERTIFICAT = 19
TEMPLATE_BUMP = 0

# Mapping formule → template(s)
FORMULE_MAP = {
    "Ligue 1":          [TEMPLATE_ORANGE],
    "Champions League":  [TEMPLATE_ORANGE],
    "FFF":              [TEMPLATE_FFF_BIENVENUE, TEMPLATE_FFF_CERTIFICAT],
    "Bump":             [TEMPLATE_BUMP],
}


# ── Parsing Excel ─────────────────────────────────────────────────────────────

def find_sheet(wb, target_date=None, sheet_name=None):
    """Trouve l'onglet correspondant à la date ou au nom donné."""
    if sheet_name:
        if sheet_name in wb.sheetnames:
            return wb[sheet_name]
        # Recherche approximative
        for name in wb.sheetnames:
            if sheet_name.lower() in name.lower():
                return wb[name]
        raise ValueError(f"Onglet '{sheet_name}' introuvable. "
                         f"Onglets disponibles : {wb.sheetnames[-10:]}")

    if target_date is None:
        target_date = datetime.now()

    # Essayer différents formats de date pour le nom d'onglet
    candidates = [
        target_date.strftime("%d.%m.%y"),           # 22.02.26
        target_date.strftime("%d.%m"),               # 22.02
        target_date.strftime("%-d.%m"),              # 2.02
    ]

    for name in wb.sheetnames:
        for candidate in candidates:
            if candidate in name:
                return wb[name]

    # Essayer avec la date en A1
    for ws in wb.worksheets:
        val = ws.cell(1, 1).value
        if isinstance(val, datetime) and val.date() == target_date.date():
            return ws

    raise ValueError(
        f"Aucun onglet trouvé pour la date {target_date.strftime('%d/%m/%Y')}. "
        f"Derniers onglets : {wb.sheetnames[-10:]}"
    )


def parse_birthdays(ws):
    """Parse les anniversaires depuis un onglet de feuille de route.

    Retourne une liste de dicts :
        {prenom, formule, horaire, nb_enfants, boisson, cadeau, options, gateau, commentaires}
    """
    birthdays = []

    # Les headers sont en ligne 5 : A=Animateur, B=Horaire, C=Formule,
    # D=Nb enfants, E=Prénom, F=Boisson, G=Cadeau, H=Options, I=Gateau, J=Commentaires
    for row in range(6, ws.max_row + 1):
        prenom = ws.cell(row, 5).value  # col E
        formule = ws.cell(row, 3).value  # col C

        if not prenom or not formule:
            continue

        prenom = str(prenom).strip()
        formule = str(formule).strip()

        if formule not in FORMULE_MAP:
            print(f"  ⚠ Formule inconnue '{formule}' pour '{prenom}' (ligne {row}), ignoré")
            continue

        entry = {
            "prenom": prenom,
            "formule": formule,
            "horaire": str(ws.cell(row, 2).value or "").strip(),
            "nb_enfants": str(ws.cell(row, 4).value or "").strip(),
            "boisson": str(ws.cell(row, 6).value or "").strip(),
            "cadeau": str(ws.cell(row, 7).value or "").strip(),
            "options": str(ws.cell(row, 8).value or "").strip(),
            "gateau": str(ws.cell(row, 9).value or "").strip(),
            "commentaires": str(ws.cell(row, 10).value or "").strip(),
            "animateur": str(ws.cell(row, 1).value or "").strip(),
            "row": row,
        }
        birthdays.append(entry)

    return birthdays


# ── Génération PPTX ──────────────────────────────────────────────────────────

def clone_slide(prs, source_index):
    """Clone un slide du PPTX en dupliquant le Part OPC."""
    source_slide = prs.slides[source_index]
    source_part = source_slide.part

    # Trouver le prochain numéro de slide disponible
    max_num = 0
    for s in prs.slides:
        m = re.search(r'slide(\d+)\.xml', str(s.part.partname))
        if m:
            max_num = max(max_num, int(m.group(1)))
    new_num = max_num + 1
    new_partname = PackURI(f'/ppt/slides/slide{new_num}.xml')

    # Copier le XML
    new_xml = copy.deepcopy(source_part._element)

    # Créer un SlidePart
    new_part = SlidePart(new_partname, source_part.content_type,
                         source_part.package, new_xml)

    # Copier les relations (images, layout)
    for rId, rel in source_part.rels.items():
        if rel.is_external:
            new_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_part.relate_to(rel.target_part, rel.reltype, rId)

    # Enregistrer le slide dans la présentation
    rId = prs.part.relate_to(new_part, RT.SLIDE)

    sldIdLst = prs.slides._sldIdLst
    existing_ids = [int(s.get('id')) for s in sldIdLst]
    new_id = max(existing_ids) + 1

    sldId = etree.SubElement(sldIdLst, qn('p:sldId'))
    sldId.set('id', str(new_id))
    sldId.set(qn('r:id'), rId)

    return prs.slides[len(prs.slides) - 1]


def delete_slide(prs, idx):
    """Supprime le slide à l'index donné."""
    sldId = prs.slides._sldIdLst[idx]
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)


def split_prenoms(prenom):
    """Sépare 'Prénom1 et Prénom2' en liste de prénoms individuels."""
    if " et " in prenom:
        return [p.strip() for p in prenom.split(" et ")]
    return [prenom]


def is_double(prenom):
    """Vérifie si c'est un double anniversaire."""
    return " et " in prenom


def set_name_on_slide(slide, template_index, prenom):
    """Remplace le prénom sur un slide selon le type de template.

    Pour les doubles anniversaires (« Prénom1 et Prénom2 ») :
    - Bump : utilise les 2 emplacements dédiés (shape[3] + shape[6])
    - Orange / FFF : affiche les 2 prénoms sur des lignes séparées
    """
    prenoms = split_prenoms(prenom)

    if template_index == TEMPLATE_BUMP:
        # Bump a 2 emplacements nom : shape[6] (visible) et shape[3] (sous la 2e image)
        if len(prenoms) >= 2:
            slide.shapes[6].text_frame.paragraphs[0].runs[0].text = prenoms[0]
            slide.shapes[3].text_frame.paragraphs[0].runs[0].text = prenoms[1]
        else:
            slide.shapes[6].text_frame.paragraphs[0].runs[0].text = prenoms[0]
            slide.shapes[3].text_frame.paragraphs[0].runs[0].text = prenoms[0]

    elif template_index == TEMPLATE_ORANGE:
        # shape[3] = prénom — pour double, on met les 2 séparés par « & »
        shape = slide.shapes[3]
        display_text = " &\n".join(prenoms) if len(prenoms) > 1 else prenoms[0]
        shape.text_frame.paragraphs[0].runs[0].text = display_text

    elif template_index == TEMPLATE_FFF_BIENVENUE:
        # shape[2] = prénom
        shape = slide.shapes[2]
        display_text = " &\n".join(prenoms) if len(prenoms) > 1 else prenoms[0]
        shape.text_frame.paragraphs[0].runs[0].text = display_text


def set_code_on_slide(slide, code):
    """Remplace le code supporter FFF sur un slide certificat."""
    # shape[1] = code
    shape = slide.shapes[1]
    shape.text_frame.paragraphs[0].runs[0].text = code


def generate_pptx(birthdays, template_path, output_path):
    """Génère le PPTX avec les affiches pour tous les anniversaires."""
    with open(template_path, 'rb') as f:
        template_bytes = f.read()

    prs = Presentation(io.BytesIO(template_bytes))
    original_count = len(prs.slides)

    slides_created = []

    for entry in birthdays:
        formule = entry["formule"]
        prenom = entry["prenom"]
        template_indices = FORMULE_MAP[formule]

        for tmpl_idx in template_indices:
            if tmpl_idx == TEMPLATE_FFF_CERTIFICAT:
                # Certificat FFF : un par enfant (chacun a son propre code)
                for individual in split_prenoms(prenom):
                    new_slide = clone_slide(prs, tmpl_idx)
                    code = entry.get("code_fff", "________")
                    set_code_on_slide(new_slide, code)
                    slides_created.append({
                        "prenom": individual,
                        "formule": formule,
                        "template": tmpl_idx,
                        "slide_index": len(prs.slides) - 1,
                    })
            else:
                # Bienvenue / poster : un seul slide avec les 2 prénoms
                new_slide = clone_slide(prs, tmpl_idx)
                set_name_on_slide(new_slide, tmpl_idx, prenom)
                slides_created.append({
                    "prenom": prenom,
                    "formule": formule,
                    "template": tmpl_idx,
                    "slide_index": len(prs.slides) - 1,
                })

    # Supprimer les slides originaux du template (indices 0 à original_count-1)
    for i in range(original_count - 1, -1, -1):
        delete_slide(prs, i)

    # Sauvegarder
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    return slides_created


def convert_to_pdf(pptx_path):
    """Convertit le PPTX en PDF via LibreOffice."""
    output_dir = os.path.dirname(pptx_path)
    try:
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf',
             '--outdir', output_dir, pptx_path],
            capture_output=True, text=True, timeout=120,
        )
        pdf_path = pptx_path.rsplit('.', 1)[0] + '.pdf'
        if os.path.exists(pdf_path):
            return pdf_path
        print(f"  ⚠ Conversion PDF échouée : {result.stderr.strip()}")
    except FileNotFoundError:
        print("  ⚠ LibreOffice non installé, pas de conversion PDF")
    except subprocess.TimeoutExpired:
        print("  ⚠ Timeout lors de la conversion PDF")
    return None


# ── Récapitulatif ─────────────────────────────────────────────────────────────

def print_recap(birthdays, sheet_title):
    """Affiche un récapitulatif des anniversaires."""
    print(f"\n{'═' * 60}")
    print(f"  ANNIVERSAIRES — {sheet_title}")
    print(f"{'═' * 60}")
    print(f"  {'Horaire':<10} {'Prénom':<20} {'Formule':<18} {'Enfants'}")
    print(f"  {'─' * 10} {'─' * 20} {'─' * 18} {'─' * 8}")

    for b in birthdays:
        print(f"  {b['horaire']:<10} {b['prenom']:<20} {b['formule']:<18} {b['nb_enfants']}")

    # Comptage par formule
    formule_counts = {}
    total_slides = 0
    for b in birthdays:
        f = b["formule"]
        formule_counts[f] = formule_counts.get(f, 0) + 1
        # Compter les slides : 1 poster par entrée + 1 certificat par enfant FFF
        n_prenoms = len(split_prenoms(b["prenom"]))
        for tmpl in FORMULE_MAP[f]:
            if tmpl == TEMPLATE_FFF_CERTIFICAT:
                total_slides += n_prenoms  # un certificat par enfant
            else:
                total_slides += 1  # un poster avec les 2 noms

    print(f"\n  Total : {len(birthdays)} anniversaires, {total_slides} slides")
    for f, count in sorted(formule_counts.items()):
        print(f"    {f}: {count}")
    print()


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Génère les affiches anniversaire UrbanSoccer")
    parser.add_argument("--date", help="Date au format DD.MM.YY (ex: 22.02.26)")
    parser.add_argument("--sheet", help="Nom exact de l'onglet Excel")
    parser.add_argument("--excel", default=EXCEL_PATH,
                        help="Chemin du fichier Excel feuille de route")
    parser.add_argument("--template", default=TEMPLATE_PATH,
                        help="Chemin du fichier PPTX template")
    parser.add_argument("--output-dir", default=OUTPUT_DIR,
                        help="Dossier de sortie")
    parser.add_argument("--no-pdf", action="store_true",
                        help="Ne pas convertir en PDF")
    args = parser.parse_args()

    # Déterminer la date cible
    target_date = None
    if args.date:
        try:
            target_date = datetime.strptime(args.date, "%d.%m.%y")
        except ValueError:
            try:
                target_date = datetime.strptime(args.date, "%d.%m.%Y")
            except ValueError:
                print(f"Format de date invalide : '{args.date}'. Utiliser DD.MM.YY")
                sys.exit(1)

    # Charger Excel
    print(f"Chargement de {os.path.basename(args.excel)}...")
    wb = openpyxl.load_workbook(args.excel, data_only=True)

    try:
        ws = find_sheet(wb, target_date, args.sheet)
    except ValueError as e:
        print(f"Erreur : {e}")
        sys.exit(1)

    sheet_title = ws.title
    print(f"Onglet trouvé : '{sheet_title}'")

    # Parser les anniversaires
    birthdays = parse_birthdays(ws)
    if not birthdays:
        print("Aucun anniversaire trouvé dans cet onglet.")
        sys.exit(0)

    print_recap(birthdays, sheet_title)

    # Générer le PPTX
    # Nom de fichier basé sur la date
    safe_name = sheet_title.replace("/", "-").replace(" ", "_")
    pptx_filename = f"Affiches_Anniversaires_{safe_name}.pptx"
    pptx_path = os.path.join(args.output_dir, pptx_filename)

    print(f"Génération du PPTX...")
    slides = generate_pptx(birthdays, args.template, pptx_path)
    print(f"  → {len(slides)} slides générés dans {pptx_path}")

    for s in slides:
        tmpl_name = {
            TEMPLATE_ORANGE: "Orange",
            TEMPLATE_FFF_BIENVENUE: "FFF Bienvenue",
            TEMPLATE_FFF_CERTIFICAT: "FFF Certificat",
            TEMPLATE_BUMP: "Bump",
        }.get(s["template"], "?")
        print(f"    Slide {s['slide_index']}: {s['prenom']} ({s['formule']} → {tmpl_name})")

    # Convertir en PDF
    if not args.no_pdf:
        print("Conversion en PDF...")
        pdf_path = convert_to_pdf(pptx_path)
        if pdf_path:
            print(f"  → {pdf_path}")

    print("\nTerminé !")


if __name__ == "__main__":
    main()
